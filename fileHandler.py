from flask import request, jsonify
from werkzeug.utils import secure_filename
import shutil
import os
import pickle
from langchain.document_loaders import PyPDFLoader, CSVLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import easyocr
from docx import Document
from pptx import Presentation
import pandas as pd
import logging


class FileHandler:
    def __init__(self, app, upload_folder, knowledgeBase_folder):
        self.app = app
        self.upload_folder = upload_folder
        self.knowledgeBase_folder = knowledgeBase_folder
        self.allowed_extensions = {'txt', 'pdf', 'png', 'jpg', 'jpeg', 'docx', 'pptx', 'xlsx', 'xls', 'csv'}
        self.file_processor = FileProcessor(knowledgeBase_folder)  # Initialize FileProcessor
        self.active_timestamps = None  # Unique and united time stamp for upload file and its related knowledge_base
        self.file_mapping = {}
        self.knowledge = None

        if not os.path.exists(self.upload_folder):
            os.makedirs(self.upload_folder)

        self.setup_routes()

    def setup_routes(self):
        self.app.route('/upload', methods=['POST'])(self.upload_file)
        self.app.route('/delete', methods=['POST'])(self.delete_file)
        self.app.route('/cleanup_uploads', methods=['GET'])(self.cleanup_uploads)

    def allowed_file(self, filename):
        return '.' in filename and filename.rsplit('.', 1)[1].lower() in self.allowed_extensions
    
    def get_category_folder(self, filename):
        extension = filename.rsplit('.', 1)[1].lower()
        if extension in {'jpg', 'jpeg', 'png'}:
            return 'images'
        elif extension == 'pdf':
            return 'pdfs'
        elif extension =='docx':
            return 'documents'
        elif extension == 'pptx':
            return 'slides'
        elif extension in {'xlsx', '.xls'}:
            return 'sheets'
        elif extension == 'csv':
            return 'csv'
        elif extension == 'txt':
            return 'text'
        else:
            return 'others'
    
    def save_file_mapping(self, file_id, timestamp_filename):
        self.file_mapping[file_id] = timestamp_filename

    def get_timestamp_filename(self, file_id):
        return self.file_mapping.get(file_id)
    
    def upload_file(self):
        if 'files' not in request.files:
            return jsonify({"error": "No file part"}), 400
        files = request.files.getlist('files')

        if not files or files[0].filename == '':
            return jsonify({'error': 'No files selected for uploading'}), 400
        
        timestamps = request.form.getlist('timestamps')
        file_ids = request.form.getlist('fileIds')

        uploaded_files = []
        for file, timestamp, file_id in zip(files, timestamps, file_ids):
            if file and self.allowed_file(file.filename):
                filename = secure_filename(file.filename)
                formatted_timestamp = timestamp.replace(':', '-').replace('.', '_')
                timestamp_filename = f"{formatted_timestamp}_{filename}"  # Name the uploaded files

                category_folder = self.get_category_folder(filename)
                category_path = os.path.join(self.upload_folder, category_folder)
                if not os.path.exists(category_path):
                    os.makedirs(category_path)
                file_path = os.path.join(category_path, timestamp_filename)
                file.save(file_path)
                uploaded_files.append(file_path)

                # Save the mapping of fileId to timestamp_filename
                self.save_file_mapping(file_id, timestamp_filename)

            # Pass the uploaded files to the FileProcessor for processing once user uploads the files
            self.file_processor.process_files(uploaded_files, timestamp_filename)
        
        return jsonify({"uploaded_files": uploaded_files}), 200

    def delete_file(self):
        try:
            # Retrieve filename and timestamp from the request
            data = request.json
            file_id = data.get('fileId')
            print('\n\nfileId:', file_id)

            if not file_id:
                return jsonify({"error": "File ID is missing"}), 400

            # Retrieve the filename associated with the fileId
            timestamp_filename = self.get_timestamp_filename(file_id)
            print('\n\ntimestamp_filename:', timestamp_filename)
            
            if not timestamp_filename:
                return jsonify({"error": "File not found"}), 404
            

            category_folder = self.get_category_folder(secure_filename(timestamp_filename))
            logging.debug(f"Attempting to delete file at path: {category_folder}")

            # Construct the full path to the file
            file_path = os.path.join(self.upload_folder, category_folder, timestamp_filename)
            logging.debug(f"Attempting to delete file at path: {file_path}")
            print(file_path)


            if os.path.exists(file_path):
                os.remove(file_path)
                # Optionally, remove associated knowledge base file
                kb_file_path = os.path.join(self.knowledgeBase_folder, f"{timestamp_filename}_knowledge.pkl")
                if os.path.exists(kb_file_path):
                    os.remove(kb_file_path)
                return jsonify({"success": True, "message": "File deleted successfully"}), 200
            else:
                return jsonify({"success": False,"error": "File not found"}), 404

        except Exception as e:
            logging.error(f"Error deleting file: {e}")
            return jsonify({"error": "An error occurred while deleting the file"}), 500
  
    def cleanup_uploads(self):
        if os.path.exists(self.upload_folder) and os.path.exists(self.knowledgeBase_folder):
            try:
                shutil.rmtree(self.upload_folder)
                logging.info(f'Deleted all contents in: {self.upload_folder}')
                os.makedirs(self.upload_folder)  # Recreate the upload folder
                shutil.rmtree(self.knowledgeBase_folder)
                logging.info(f'Deleted all contents in: {self.upload_folder}')
                os.makedirs(self.knowledgeBase_folder)  # Recreate the knowledge_base folder
                return jsonify({'status': 'success', 'message': 'Uploads cleaned up'}), 200
            except Exception as e:
                logging.error(f'Failed to clean up {self.upload_folder}. Reason: {e}')
        else:
            logging.warning(f"The directory {self.upload_folder} or {self.knowledgeBase_folder} does not exist")

class FileProcessor:
    def __init__(self, knowledgeBase_folder):
        self.knowledgeBase_folder = knowledgeBase_folder

    def process_files(self, file_paths, timestamp_filename):  # Get the unique time stamp from class FileHandler
        for file_path in file_paths:
            category = self.get_category(file_path)
            if category:
                self.process_file(file_path, category, timestamp_filename)

    def get_category(self, file_path):
        if 'images' in file_path:
            return 'images'
        elif 'pdfs' in file_path:
            return 'pdfs'
        elif 'documents' in file_path:
            return 'documents'
        elif 'sheets' in file_path:
            return 'sheets'
        elif 'text' in file_path:
            return 'text'
        else:
            return None

    def process_file(self, file_path, category, timestamp_filename):
        if category == 'images':
            self.process_image(file_path, timestamp_filename)
        elif category == 'pdfs':
            self.process_pdf(file_path, timestamp_filename)
        elif category == 'documents':
            self.process_document(file_path, timestamp_filename)
        elif category == 'slides':
            self.process_slide(file_path, timestamp_filename)
        elif category == 'sheets':
            self.process_sheet(file_path, timestamp_filename)
        elif category == 'csv':
            self.process_csv(file_path, timestamp_filename)
        elif category == 'text':
            self.process_text(file_path, timestamp_filename)
        else:
            print(f"Unhandled file category: {category}")

    def process_image(self, file_path, timestamp_filename):
        reader = easyocr.Reader(['en'])
        try:
            image_data = reader.readtext(file_path)
            image_text_contents = [item[1] for item in image_data]
            print(f"Processed image: {file_path}")
            self.save_to_file(image_text_contents, timestamp_filename)
        except Exception as e:
            print(f"Failed to process image file {file_path}. Error: {e}")

    def process_pdf(self, file_path, timestamp_filename):
        loader = PyPDFLoader(file_path)
        try:
            pdf_data = loader.load()
            print(f"Processed PDF: {file_path}")
            text_splitter = RecursiveCharacterTextSplitter(chunk_size = 300,chunk_overlap = 50,)
            contents = text_splitter.split_documents(pdf_data)
            self.save_to_file(contents, timestamp_filename)
        except Exception as e:
            print(f"Failed to process PDF file {file_path}. Error: {e}")

    def process_document(self, file_path, timestamp_filename):
        text = []
        doc_data = Document(file_path)
        try:
            for paragraph in doc_data.paragraphs:
                text.append(paragraph.text)
            self.save_to_file(text, timestamp_filename)
        except Exception as e:
            print(f"Failed to process document file {file_path}. Error: {e}")

    def process_slide(self, file_path, timestamp_filename):
        text = []
        slide_data = Presentation(file_path)
        try:
            for slide in slide_data.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            self.save_to_file(text, timestamp_filename)
        except Exception as e:
            print(f"Failed to process document file {file_path}. Error: {e}")

    def process_sheet(self, file_path, timestamp_filename):
        try:
            text_contents = []
            sheet_data = pd.read_excel(file_path, sheet_name=None)  # read all sheets
            
            if isinstance(sheet_data, dict):
                for sheet_name, data in sheet_data.items():
                    text_contents.append(f"Sheet: {sheet_name}\n{data.to_string()}")
            else:
                text_contents.append(sheet_data.to_string())
            
            print(f"Processd sheet file: {file_path}")
            self.save_to_file(text_contents, timestamp_filename)

        except Exception as e:
            print(f"Failed to process sheet file {file_path}. Error: {e}")

    def process_csv(self, file_path, timestamp_filename):
        loader = CSVLoader(file_path)
        try:
            csv_data = loader.load()
            print(f"Processed CSV: {file_path}")
            text_splitter = RecursiveCharacterTextSplitter(chunk_size = 300,chunk_overlap = 50,)
            contents = text_splitter.split_documents(csv_data)
            self.save_to_file(contents, timestamp_filename)
        except Exception as e:
            print(f"Failed to process CSV file {file_path}. Error: {e}")

    def process_text(self, file_path, timestamp_filename):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text_data = file.read()
                print(f"Processd text: {file_path}")
                self.save_to_file([text_data], timestamp_filename)
        except Exception as e:
            print(f"Failed to process text file {file_path}. Error: {e}")

    def save_to_file(self, contents, timestamp_filename):
        output_file = os.path.join(self.knowledgeBase_folder, f'{timestamp_filename}_knowledge.pkl')
        try:
            with open(output_file, "wb") as file:
                pickle.dump(contents, file)
            print(f"Text contents saved to {output_file}")
        except Exception as e:
            print(f"Failed to save text contents to file. Error: {e}")